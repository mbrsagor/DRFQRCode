from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "PowurHouz 2.0: Scaling for the Future"
    subtitle.text = "Transitioning to a Modern, Automated, Production-Grade Architecture\n\nPresented by: GDS Team"

    # Define content slides
    slides_data = [
        {
            "title": "The Catalyst for Change: Why Rebuild Now?",
            "bullets": [
                "The Legacy Foundation: Our MVP validated the business model over the last two years, but data volume is outgrowing the architecture.",
                "The Production Bottleneck: Monolithic design limits rapid feature deployment and scaling for complex engine and alternator data.",
                "The Automation Imperative: Manual workflows must evolve into real-time automation to stay competitive.",
                "Speed to Market: Decoupling the frontend and backend allows our teams to build and ship independently."
            ]
        },
        {
            "title": "Current State vs. Limitations: The Monolith",
            "bullets": [
                "Tightly Coupled Architecture: UI and business logic are intertwined, making simple dashboard updates unnecessarily complex.",
                "Synchronous Processing: Heavy data queries block the system, resulting in slower load times.",
                "Static User Experience: Traditional Bootstrap templates lack the dynamic, instant-feedback feel of modern SaaS platforms.",
                "Scaling Challenges: Impossible to scale specific high-traffic services independently without scaling the entire application."
            ]
        },
        {
            "title": "The Vision for PowurHouz 2.0",
            "bullets": [
                "API-First Design: A robust central API serving the web dashboard, future mobile apps, and third-party integrations.",
                "Asynchronous Power: Background workers handling the heavy lifting without freezing the user interface.",
                "Reactive Frontend: A highly responsive, component-based UI that feels like native desktop software.",
                "AI-Ready: Infrastructure designed from day one to ingest and deploy machine learning models."
            ]
        },
        {
            "title": "The Modern Backend & AI Integration",
            "bullets": [
                "Core API (Python, Django, DRF): Providing clean, secure, and standardized JSON endpoints.",
                "Async Engine (Redis, Celery, RabbitMQ): Enterprise message brokering for background tasks, data syncing, and scheduled reports.",
                "AI Feature 1: Predictive Inventory: AI models forecasting engine/radiator requirements based on historical seasonal trends.",
                "AI Feature 2: Smart Lead Scoring: Automated categorization and routing of incoming leads.",
                "AI Feature 3: Anomaly Detection: Automated flagging of incorrect system model and serial number entries."
            ]
        },
        {
            "title": "The Modern Frontend Experience",
            "bullets": [
                "React.js: Transitioning to reusable, testable UI components for faster development.",
                "Material-UI (MUI): Utilizing an enterprise-grade component library for a consistent, professional design language.",
                "Modern Admin Template: Accelerating development with a premium, React-ready foundation for the dashboard.",
                "Real-Time State Management: Instantaneous UI updates without full page reloads, providing a seamless user experience."
            ]
        },
        {
            "title": "Infrastructure & Deployment Strategy",
            "bullets": [
                "Containerization (Docker): Packaging the entire stack into isolated containers to guarantee consistency.",
                "CI/CD Pipelines: Automated testing and deployment workflows (via GitHub Actions) for safe, rapid updates.",
                "Cloud Infrastructure: Deploying on scalable cloud architecture (e.g., AWS EC2 and S3) to handle traffic spikes."
            ]
        },
        {
            "title": "The Roadmap to Go-Live",
            "bullets": [
                "Phase 1: API Foundation. Build the Django REST API alongside the legacy system; connect to the existing database.",
                "Phase 2: UI Rebuild. Develop the React/MUI frontend simultaneously, consuming the new API endpoints.",
                "Phase 3: Automation. Implement Celery, Redis, and RabbitMQ for background tasks and AI features.",
                "Phase 4: Containerization. Dockerize the environments and establish the CI/CD pipeline.",
                "Phase 5: Launch. Switch user traffic to the modern stack with zero downtime."
            ]
        }
    ]

    # Add content slides
    bullet_slide_layout = prs.slide_layouts[1]
    for slide_data in slides_data:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = slide_data["title"]
        
        tf = body_shape.text_frame
        tf.text = slide_data["bullets"][0]
        
        for bullet in slide_data["bullets"][1:]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

    # Save presentation
    prs.save('PowurHouz_2.0_Presentation.pptx')
    print("Presentation saved as 'PowurHouz_2.0_Presentation.pptx'")

if __name__ == "__main__":
    create_presentation()

